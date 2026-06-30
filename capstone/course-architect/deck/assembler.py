"""The assembler tool: a DeckSpec becomes a real .pptx via python-pptx.

This is pure layout. It reads the content the agent produced and lays it out on
16:9 slides with the chosen theme. It adds NO content of its own — if a bullet
isn't in the spec, it won't appear in the deck. (That property is what the coverage
evaluator checks.)
"""

from __future__ import annotations

from typing import List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .spec import DeckSpec, Section, Slide
from .theme import Theme, get_theme

_EMU_W, _EMU_H = Inches(13.333), Inches(7.5)  # 16:9


def _rgb(hex6: str) -> RGBColor:
    return RGBColor(int(hex6[0:2], 16), int(hex6[2:4], 16), int(hex6[4:6], 16))


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank


def _set_bg(slide, hex6: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(hex6)


def _accent_band(slide, theme: Theme, height_in: float = 0.18, top_in: float = 0.0) -> None:
    box = slide.shapes.add_textbox(0, Inches(top_in), _EMU_W, Inches(height_in))
    box.fill.solid(); box.fill.fore_color.rgb = _rgb(theme.accent)
    box.line.fill.background()


def _textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _style_run(run, size, color_hex, font, bold=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color_hex)


def _footer(slide, theme: Theme, text: str) -> None:
    tf = _textbox(slide, Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    _style_run(run, 10, theme.muted, theme.font_body)


# --------------------------------------------------------------------------- #

def _title_slide(prs, deck: DeckSpec, theme: Theme) -> None:
    slide = _blank(prs)
    _set_bg(slide, theme.bg)
    _accent_band(slide, theme, height_in=2.6, top_in=2.45)  # central band
    tf = _textbox(slide, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = deck.title
    _style_run(r, 40, "FFFFFF", theme.font_title, bold=True)
    if deck.subtitle or deck.audience:
        sub = _textbox(slide, Inches(0.9), Inches(4.15), Inches(11.5), Inches(0.8))
        p2 = sub.paragraphs[0]
        r2 = p2.add_run()
        r2.text = deck.subtitle or f"For: {deck.audience}"
        _style_run(r2, 18, "FFFFFF", theme.font_body)


def _section_slide(prs, section: Section, theme: Theme, idx: int) -> None:
    slide = _blank(prs)
    _set_bg(slide, theme.band)
    tf = _textbox(slide, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"Module {idx}"
    _style_run(r, 16, theme.accent, theme.font_body, bold=True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = section.module_title
    _style_run(r2, 34, "FFFFFF", theme.font_title, bold=True)
    if section.objective:
        p3 = tf.add_paragraph()
        r3 = p3.add_run(); r3.text = section.objective
        _style_run(r3, 16, "E5E7EB", theme.font_body)


def _content_slide(prs, slide_spec: Slide, theme: Theme, footer: str) -> None:
    slide = _blank(prs)
    _set_bg(slide, theme.bg)
    _accent_band(slide, theme)
    # Title
    ttf = _textbox(slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.0))
    tp = ttf.paragraphs[0]
    tr = tp.add_run(); tr.text = slide_spec.title
    _style_run(tr, 28, theme.title, theme.font_title, bold=True)
    # Bullets
    btf = _textbox(slide, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.0))
    first = True
    for bullet in slide_spec.bullets:
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        marker = p.add_run(); marker.text = "•  "
        _style_run(marker, 18, theme.accent, theme.font_body, bold=True)
        run = p.add_run(); run.text = bullet
        _style_run(run, 18, theme.body, theme.font_body)
    if not slide_spec.bullets:
        btf.paragraphs[0].add_run().text = ""
    _footer(slide, theme, footer)
    # Speaker notes
    if slide_spec.notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.notes


def _summary_slide(prs, deck: DeckSpec, theme: Theme) -> None:
    slide = _blank(prs)
    _set_bg(slide, theme.bg)
    _accent_band(slide, theme)
    ttf = _textbox(slide, Inches(0.6), Inches(0.45), Inches(12.1), Inches(1.0))
    tr = ttf.paragraphs[0].add_run(); tr.text = "Course summary"
    _style_run(tr, 28, theme.title, theme.font_title, bold=True)
    btf = _textbox(slide, Inches(0.7), Inches(1.7), Inches(11.9), Inches(5.0))
    first = True
    for i, sec in enumerate(deck.sections, 1):
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        marker = p.add_run(); marker.text = f"{i}.  "
        _style_run(marker, 18, theme.accent, theme.font_body, bold=True)
        run = p.add_run(); run.text = sec.module_title + (f" — {sec.objective}" if sec.objective else "")
        _style_run(run, 18, theme.body, theme.font_body)


def assemble_pptx(deck: DeckSpec, out_path: str, theme: Theme | None = None) -> str:
    """Build the .pptx and return its path."""
    th = theme or get_theme(deck.theme)
    prs = Presentation()
    prs.slide_width, prs.slide_height = _EMU_W, _EMU_H

    _title_slide(prs, deck, th)
    n = deck.content_slide_count()
    page = 0
    for i, section in enumerate(deck.sections, 1):
        _section_slide(prs, section, th, i)
        for sl in section.slides:
            page += 1
            _content_slide(prs, sl, th, f"{deck.title}   ·   {page}/{n}")
    _summary_slide(prs, deck, th)

    prs.save(out_path)
    return out_path
